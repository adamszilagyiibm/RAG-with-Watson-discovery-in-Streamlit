
# Important: hardcoding the API key in Python code is not a best practice. We are using
# this approach for the ease of demo setup. In a production application these variables
# can be stored in anx00 .env or a properties file

# For reading credentials from the .env file
import os
from dotenv import load_dotenv
import chromadb
from pptx import Presentation
from langchain.text_splitter import CharacterTextSplitter
from langchain_community.document_loaders import TextLoader
from langchain_community.document_loaders import PyPDFLoader
from chromadb.utils import embedding_functions
from io import BytesIO
from pptx.util import Inches, Pt
from ibm_watson import DiscoveryV2
from ibm_cloud_sdk_core.authenticators import IAMAuthenticator
import json

# WML python SDK
from ibm_watson_machine_learning.foundation_models import Model
from ibm_watson_machine_learning.metanames import GenTextParamsMetaNames as GenParams
from ibm_watson_machine_learning.foundation_models.utils.enums import ModelTypes, DecodingMethods

# URL of the hosted LLMs is hardcoded because at this time all LLMs share the same endpoint
FILE_TYPE_TXT = "txt"
FILE_TYPE_PDF = "pdf"
FILE_TYPE_PPT = "ppt"
global_response = ""
ppt_response = ""
file_type = ""
default_ef = embedding_functions.DefaultEmbeddingFunction()



# Replace with your watsonx project id (look up in the project Manage tab)
watsonx_project_id = ""
# Replace with your IBM Cloud key
#api_key = ""
url = ""
watsonDiscoveryAPI_key = ""
watsonDiscoveryURL = ""
watsonDiscoveryproject_id=""
WatsonDisvoerycollection_ids=""
c_api_key =""
def get_credentials():

    load_dotenv()
    # Update the global variables that will be used for authentication in another function
    globals()["c_api_key"] = os.getenv("cloud_api_key", None)
    globals()["watsonx_project_id"] = os.getenv("project_id", None)
    globals()["watsonDiscoveryAPI_key"] = os.getenv("watsonDiscoveryAPIkey", None)
    globals()["watsonDiscoveryURL"] = os.getenv("watsonDiscoveryURL", None)
    globals()["watsonDiscoveryproject_id"] = os.getenv("watsonDiscoveryproject_id", None)
    globals()["WatsonDisvoerycollection_ids"] = os.getenv("WatsonDisvoerycollection_ids", None)

    print(c_api_key)
# The get_model function creates an LLM model object with the specified parameters

def get_model(model_type,max_tokens,min_tokens,decoding,temperature):

    generate_params = {
        GenParams.MAX_NEW_TOKENS: max_tokens,
        GenParams.MIN_NEW_TOKENS: min_tokens,
        GenParams.DECODING_METHOD: decoding,
        GenParams.TEMPERATURE: temperature
    }

    model = Model(
        model_id=model_type,
        params=generate_params,
        credentials={
            "apikey": c_api_key,
            "url": url
        },
        project_id=watsonx_project_id
        )

    return model


import streamlit as st
import fitz  # PyMuPDF
 #import use_case_RAG

def main():

    get_credentials()

    # Declare variables
    file_path = ""
    collection_name = ""
    file_type = ""
    FILE_TYPE_TXT = "txt"
    FILE_TYPE_PDF = "pdf"
    FILE_TYPE_PPT = "ppt"
    # Use the full page instead of a narrow central column
    st.set_page_config(layout="wide")

    # Streamlit app title
    st.title("Demo szöveg generálás watson discoveryből data and ai-ra")
    # Write bold text
    # UI component to enter the question
    question = st.text_area('Kérdés',height=100)

    if 'global_response' not in st.session_state:
        st.session_state.global_response = ""
    button_clicked = st.button("Válaszold meg a kérdést")

    st.subheader("Válasz:")

    # Invoke the LLM when the button is clicked

    if button_clicked:
        def answer_questions_from_doc(c_api_key,watsonx_project_id,watsonDiscoveryAPI_key,watsonDiscoveryURL,watsonDiscoveryproject_id,WatsonDisvoerycollection_ids):
            # print("watson discovery api key" + watsonDiscoveryAPI_key)
            # authenticator = IAMAuthenticator('your Watson Discovery API key')
            authenticator = IAMAuthenticator(watsonDiscoveryAPI_key)
            discovery = DiscoveryV2(
            version='2020-08-30',
            authenticator=authenticator
            )
            discovery.set_service_url(watsonDiscoveryURL)

            response = discovery.query(
                project_id = watsonDiscoveryproject_id,
                collection_ids = [WatsonDisvoerycollection_ids],
                passages = {'enabled': True, 
                            'max_per_document': 20,
                            'find_answers': True},
                natural_language_query = question
            ).get_result()

            with open('data.json', 'w') as f:
                json.dump(response, f)
            # inspecting the key fields in the WD output

            response.keys()

            # only one relevant document (because one document in the collection)
            len(response['results'])

            # removing some tags
            passages = response['results'][0]['document_passages']
            passages = [p['passage_text'].replace('<em>', '').replace('</em>', '').replace('\n','') for p in passages]
            passages

            # concatenating passages
            context = '\n '.join(passages)
            context
            prompt = \
            "Kérlek válaszold meg a kérdést a következő dokumentumból. \
            Ha kérdés megválaszolhatatlan, mond 'megválaszolhatatlan'. \
            Ha megválaszoltad a kérdést, ne mond 'megválaszolhatatlan'. \
            Ne Tegyél bele információt, ami nem releváns a kérdésre. \
            Ne válaszolj meg más kérdéseket \
            Legyél biztos benne, hogy a felhasznált nyelv magyar legyen.'\
            Ne használj ismétlődést' ha a válaszban angol szöveg van fordítsd le magyar nyelvre\
            Question:" + question +  context + "a válasz csak és kizárólag magyarul legyen"\

            complete_prompt = prompt + question

            #print("----------------------------------------------------------------------------------------------------")
            #print("*** Prompt:" + prompt + "***")
            #print("----------------------------------------------------------------------------------------------------")

    # URL of the hosted LLMs is hardcoded because at this time all LLMs share the same endpoint
            
            def get_model(model_type,max_tokens,min_tokens,decoding,temperature):#, repetition_penalty):

                generate_params = {
                GenParams.MAX_NEW_TOKENS: max_tokens,
                GenParams.MIN_NEW_TOKENS: min_tokens,
                GenParams.DECODING_METHOD: decoding,
                GenParams.TEMPERATURE: temperature,
                }
    
                model = Model(
                model_id=model_type,
                params=generate_params,
                credentials={
                    "apikey": c_api_key,
                    "url": url
                },
                project_id= watsonx_project_id
                )

                return model
            # api_key =""
            url = "https://us-south.ml.cloud.ibm.com"
            # Replace with your watsonx project id (look up in the project Manage tab)
            watsonx_project_id = watsonx_project_id
            # Replace with your IBM Cloud key
            api_key = c_api_key  

            model_type = "mistralai/mixtral-8x7b-instruct-v01"
            max_tokens = 800
            min_tokens = 100
            decoding = DecodingMethods.GREEDY
            temperature = 0.7
                
            model = get_model(model_type, max_tokens, min_tokens, decoding, temperature)
            #model_type = "meta-llama/llama-2-70b-chat"
            # model_type = "google/flan-t5-xxl"
            # model_type = "ibm/granite-13b-chat-v1"
            # model_type = "ibm/granite-13b-instruct-v1"
            # model_id = "ibm/mpt-7b-instruct2"
            #max_tokens = 800
            #min_tokens = 50
            #decoding = DecodingMethods.GREEDY
            #temperature = 0.7

            response = model.generate(complete_prompt)
            response_text = response['results'][0]['generated_text']
            #response_text = response['generated_text']
            # print model respons
            print("Válasz az LLM-től:" + response_text)
            return response_text
            
        
        response=answer_questions_from_doc(c_api_key,watsonx_project_id,watsonDiscoveryAPI_key,watsonDiscoveryURL,watsonDiscoveryproject_id,WatsonDisvoerycollection_ids)
        print("Válasz az LLM-től:" + response)
        
        response.replace(response[:2], '')  

        st.write("Válasz az LLM-től:" + response)
        st.session_state.global_response=response
        
        #st.write(st.session_state.global_response)
    global_response=st.session_state.global_response
    print(global_response)

        
    prs = Presentation()   
    #title_slide_layout = prs.slide_layouts[0]
     #slide = prs.slides.add_slide(title_slide_layout)
     #title = slide.shapes.title
     #subtitle = slide.placeholders[1]

    SLD_LAYOUT_TITLE_AND_CONTENT = 1
    slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
    slide = prs.slides.add_slide(slide_layout)
    #for shape in slide.shapes:
     #       if not shape.has_text_frame:
      #          continue          
     #text_frame.clear()
     #p = text_frame.paragraphs[0]
     #run = p.add_run()
     #run.text = global_response
     #t ext_frame = slide.shapes

    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Watsonx válasz'
    tf = body_shape.text_frame
    tf.text = global_response

    # save the output into binary form
    binary_output = BytesIO()
    prs.save(binary_output) 

    st.download_button(label = 'Download pptx',
        data = binary_output.getvalue(),
        file_name = 'my_power.pptx')

if __name__ == "__main__":
    main()