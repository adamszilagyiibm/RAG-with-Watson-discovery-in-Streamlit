
# Important: hardcoding the API key in Python code is not a best practice. We are using
# this approach for the ease of demo setup. In a production application these variables
# can be stored in anx00 .env or a properties file

# For reading credentials from the .env file
import os
from dotenv import load_dotenv
import chromadb
from pptx import Presentation
from langchain.text_splitter import CharacterTextSplitter
from langchain_community.document_loaders import TextLoader
from langchain_community.document_loaders import PyPDFLoader
from chromadb.utils import embedding_functions
from io import BytesIO
from pptx.util import Inches, Pt
from ibm_watson import DiscoveryV2
from ibm_cloud_sdk_core.authenticators import IAMAuthenticator

# WML python SDK
from ibm_watson_machine_learning.foundation_models import Model
from ibm_watson_machine_learning.metanames import GenTextParamsMetaNames as GenParams
from ibm_watson_machine_learning.foundation_models.utils.enums import ModelTypes, DecodingMethods

# URL of the hosted LLMs is hardcoded because at this time all LLMs share the same endpoint
FILE_TYPE_TXT = "txt"
FILE_TYPE_PDF = "pdf"
FILE_TYPE_PPT = "ppt"
global_response = ""
ppt_response = ""
file_type = ""
default_ef = embedding_functions.DefaultEmbeddingFunction()

# Replace with your watsonx project id (look up in the project Manage tab)
watsonx_project_id = ""
# Replace with your IBM Cloud key
api_key = ""
url = ""
watsonDiscoveryAPIkey = ""
watsonDiscoveryURL = ""
watsonDiscoveryproject_id=""
WatsonDisvoerycollection_ids=""
def get_credentials():

    load_dotenv()
    # Update the global variables that will be used for authentication in another function
    globals()["api_key"] = os.getenv("api_key", None)
    globals()["watsonx_project_id"] = os.getenv("project_id", None)
    globals()["watsonDiscoveryAPIkey"] = os.getenv("watsonDiscoveryAPIkey", None)
    globals()["watsonDiscoveryURL"] = os.getenv("watsonDiscoveryURL", None)

# The get_model function creates an LLM model object with the specified parameters

# authenticator = IAMAuthenticator('your Watson Discovery API key')
authenticator = IAMAuthenticator(watsonDiscoveryAPIkey)
discovery = DiscoveryV2(
    version='2020-08-30',
    authenticator=authenticator
)

#discovery.set_service_url('Your Discovery URL')
discovery.set_service_url(watsonDiscoveryURL)


def get_model(model_type,max_tokens,min_tokens,decoding,temperature):

    generate_params = {
        GenParams.MAX_NEW_TOKENS: max_tokens,
        GenParams.MIN_NEW_TOKENS: min_tokens,
        GenParams.DECODING_METHOD: decoding,
        GenParams.TEMPERATURE: temperature
    }

    model = Model(
        model_id=model_type,
        params=generate_params,
        credentials={
            "apikey": api_key,
            "url": url
        },
        project_id=watsonx_project_id
        )

    return model

def create_embedding(file_path,file_type,collection_name):

    if file_type == FILE_TYPE_TXT:
        loader = TextLoader(file_path,encoding="1252")
        documents = loader.load()
    elif file_type == FILE_TYPE_PDF:
        loader = PyPDFLoader(file_path)
        documents = loader.load()

    text_splitter = CharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    texts = text_splitter.split_documents(documents)

    print(type(texts))

    # Load chunks into chromadb
    client = chromadb.Client()
    collection = client.get_or_create_collection(collection_name,embedding_function=default_ef)
    collection.upsert(
        documents=[doc.page_content for doc in texts],
        ids=[str(i) for i in range(len(texts))],  # unique for each doc
    )

    return collection

def create_prompt(file_path, file_type, question, collection_name):

    # Create embeddings for the text file
    collection = create_embedding(file_path,file_type,collection_name)

    # Query relevant information
    # You can try retrieving different number of chunks (n_results)
    relevant_chunks = collection.query(
        query_texts=[question],
        n_results=5,
    )

    context = "\n\n\n".join(relevant_chunks["documents"][0])
    # Please note that this is a generic format. You can change this format to be specific to llama
    prompt =  (f"{context}\n\nKérlek válaszold meg a kérdést magyarul és csak a szöveg alapján, ne legyen benne ismétlés"
              + f"  Kérlek válaszold meg a kérdést a következő dokumentumból. \
                    Ha kérdés megválaszolhatatlan, mond 'megválaszolhatatlan'. \
                    Ha megválaszoltad a kérdést, ne mond 'megválaszolhatatlan'. \
                    Ne Tegyél bele információt, ami nem releváns a kérdésre. \
                    Ne válaszolj meg más kérdéseket \
                    Legyél biztos benne, hogy a felhasznált nyelv magyar legyen.'\
                    Ne használj ismétlődést' "
              + f"  ha nem tudod megválaszolni, mond \"megválaszolhatatlan\"."
              + f"{question}")
    

    return prompt

import streamlit as st
import fitz  # PyMuPDF
 #import use_case_RAG

def main():

    get_credentials()

    # Declare variables
    file_path = ""
    collection_name = ""
    file_type = ""
    FILE_TYPE_TXT = "txt"
    FILE_TYPE_PDF = "pdf"
    FILE_TYPE_PPT = "ppt"
    # Use the full page instead of a narrow central column
    st.set_page_config(layout="wide")

    # Streamlit app title
    st.title("Demo of RAG with files and in-memory chromadb")

    # Write bold text
    st.markdown('<font color="blue"><b><i>Please upload one file type at a time and delete the file if switching types of files.</i></b></font>', unsafe_allow_html=True)

    # UI component for uploading a PDF file
    pdf_file = st.file_uploader("Upload a PDF File", type=["pdf"])

    # UI component for uploading a TXT file
    txt_file = st.file_uploader("Upload a TXT File", type=["txt"])


    # Check if a PDF file is uploaded
    if pdf_file:

        # Load the PDF content
        pdf_data = pdf_file.read()
        pdf_doc = fitz.open(stream=pdf_data, filetype="pdf")

        # Used for debugging
        print("Name of the uploaded pdf_file:" + pdf_file.name)

        # For writing file to disk
        full_file_name = pdf_file.name

        # Generate a unique collection name that follows chhoma's standards for colleciton names
        collection_name = pdf_file.name.lower()
        # Remove the .pdf
        collection_name= "pdf_" + collection_name[:-4]
        # For debugging
        print("collection_name: " + collection_name)

        # Save the pdf file in the current directory, it will be used by the module that loads data into chromadb
        pdf_doc.save(full_file_name)

        # Parameters to invoke the RAG module
        file_path = full_file_name
        file_type = FILE_TYPE_PDF

        # Close the PDF document
        pdf_doc.close()

    elif txt_file:

        # Read and display the content of the uploaded text file
        file_content = txt_file.read()
        # Convert bytes to string
        file_content = file_content.decode('latin-1')

        # Generate a unique collection name that follows chhoma's standards for colleciton names
        collection_name = txt_file.name.lower()
        # Remove the .txt
        collection_name = "txt_" + collection_name[:-4]
        # For debugging
        print("collection_name: " + collection_name)

        # Parameters to invoke the RAG module
        full_file_name = txt_file.name
        file_path = full_file_name
        file_type = FILE_TYPE_TXT

        # Define the path to the output text file
        # Write the content to the output text file
        with open(full_file_name, 'w', encoding='utf-8') as output_file:
            output_file.write(file_content)


    # UI component to enter the question
    question = st.text_area('Kérdés',height=100)

    if 'global_response' not in st.session_state:
        st.session_state.global_response = ""
    button_clicked = st.button("Válaszold meg a kérdést")

    st.subheader("Válasz:")

    # Invoke the LLM when the button is clicked

    if button_clicked:
       # try:
        def answer_questions_from_doc(api_key, watsonx_project_id, file_path, FILE_TYPE_PDF,question,collection_name):
        # Update the global variable
        # Specify model parameters
            print("model inicial")
            model_type = "mistralai/mixtral-8x7b-instruct-v01"
            max_tokens = 800
            min_tokens = 100
            decoding = DecodingMethods.GREEDY
            temperature = 0.7
            # Get the watsonx model
            model = get_model(model_type, max_tokens, min_tokens, decoding, temperature)
            # Get the prompt
            print("prompt generálás")
            complete_prompt = create_prompt(file_path, file_type, question, collection_name)
            print("prompt kész")
            generated_response = model.generate(prompt=complete_prompt)
            print("modell válaszolt")
            response_text = generated_response['results'][0]['generated_text']
            print("válasz generálás")
            return response_text
        response=answer_questions_from_doc(api_key,watsonx_project_id,file_path,file_type,question,collection_name)
        print("doc")
        print("Válasz az LLM-től:" + response)
        st.write(response)
        print("global variable mentés")
        print("global variable done ")
        st.session_state.global_response=response
    
    #st.write(st.session_state.global_response)
    global_response=st.session_state.global_response
    print(global_response)
        # UI component for editing a PPT file
     #   finally:
     # ppt_file = st.file_uploader("Upload a PPT File", type=["ppt"])
    # Load the PDF content
    
    #uploaded_files = st.file_uploader("Choose a PPTX file",type=["pptx"],
    #                              accept_multiple_files=True)
                                  
    #def upload():
    #    if not uploaded_files:
    #        st.text('Load pptx file/s first')
    #    else:
    #        for uploaded_file in uploaded_files:
    #            prs = Presentation(uploaded_file)
    #            prs2 =  Presentation(uploaded_file)
    #            SLD_LAYOUT_TITLE_AND_CONTENT = 1
    #            slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
    #            slide = prs.slides.add_slide(slide_layout)
    #prs2 = st.write(st.session_state.result)
    #st.button("Parse text out of pptx file/s", on_click=upload, disabled=not uploaded_files)
    
    prs = Presentation()   
    #title_slide_layout = prs.slide_layouts[0]
     #slide = prs.slides.add_slide(title_slide_layout)
     #title = slide.shapes.title
     #subtitle = slide.placeholders[1]

    SLD_LAYOUT_TITLE_AND_CONTENT = 1
    slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
    slide = prs.slides.add_slide(slide_layout)
    for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            text_frame.text = global_response
    # save the output into binary form
    binary_output = BytesIO()
    prs.save(binary_output) 

    st.download_button(label = 'Download pptx',
                    data = binary_output.getvalue(),
                    file_name = 'my_power.pptx')

if __name__ == "__main__":
    main()